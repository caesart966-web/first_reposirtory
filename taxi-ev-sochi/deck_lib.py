"""Оформительский слой под корпоративный шаблон GRP.

Все константы вытащены из самого шаблона (`assets/template.pptx`):
рабочая область, цветовая схема, размеры шрифтов, макеты.
"""
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# ------------------------------------------------------------------ геометрия
SLIDE_W, SLIDE_H = 13.333, 7.5
LEFT, RIGHT = 0.37, 12.97          # границы рабочей области
TOP, BOTTOM = 1.31, 6.82
CW = RIGHT - LEFT                  # 12.60
CH = BOTTOM - TOP                  # 5.51

# -------------------------------------------------------------------- палитра
TEXT = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x7F, 0x7F, 0x7F)          # размерность в названиях графиков
GRAY_LT = RGBColor(0xD9, 0xD9, 0xD9)
ACCENT = RGBColor(0x16, 0x05, 0xFC)        # акцентный текст
ICON = RGBColor(0x10, 0x04, 0xBD)          # иконки, булиты
NEG = RGBColor(0xC0, 0x00, 0x00)
POS = RGBColor(0x00, 0x96, 0x44)
LIGHT = RGBColor(0xE4, 0xF8, 0xFE)         # вспомогательная заливка фигур
AQUA = RGBColor(0xCE, 0xEB, 0xF2)
MINT = RGBColor(0xD3, 0xEF, 0xE9)
NEG_LT = RGBColor(0xF7, 0xDD, 0xDD)
POS_LT = RGBColor(0xDD, 0xF0, 0xE5)
# основная гамма графиков: от тёмного к светлому
C1 = RGBColor(0x04, 0x3D, 0x4A)
C2 = RGBColor(0x07, 0x73, 0x8D)
C3 = RGBColor(0x0A, 0xC2, 0xF1)
C4 = RGBColor(0x86, 0xE1, 0xF9)
C5 = RGBColor(0xD0, 0xF4, 0xFD)
GAMMA = [C1, C2, C3, C4, C5]

FONT = 'Calibri'


# ------------------------------------------------------------------ утилиты
def _fmt(run, size=11, bold=False, color=TEXT, italic=False, font=FONT):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    run.font.color.rgb = color
    # кириллица должна брать тот же шрифт, что и латиница
    rPr = run._r.get_or_add_rPr()
    for tag in ('a:ea', 'a:cs'):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set('typeface', font)


def fill_tf(tf, lines, size=11, bold=False, color=TEXT, align=PP_ALIGN.LEFT,
            space_after=0, line_spacing=1.0, bullet=False, italic=False):
    """Заполняет text_frame списком строк.

    Элемент списка — str или список кортежей (текст, kwargs) для смешанного
    форматирования внутри абзаца.
    """
    tf.word_wrap = True
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        chunks = line if isinstance(line, list) else [(line, {})]
        if bullet:
            _set_bullet(p)
        for text, kw in chunks:
            r = p.add_run()
            r.text = text
            _fmt(r, **{'size': size, 'bold': bold, 'color': color,
                       'italic': italic, **kw})
    return tf


def _set_bullet(p):
    pPr = p._p.get_or_add_pPr()
    pPr.set('marL', '114300')
    pPr.set('indent', '-114300')
    for tag, attrs in (('a:buClr', None), ('a:buFont', {'typeface': 'Arial'}),
                       ('a:buChar', {'char': '▪'})):
        el = pPr.makeelement(qn(tag), attrs or {})
        if tag == 'a:buClr':
            clr = el.makeelement(qn('a:srgbClr'), {'val': '1004BD'})
            el.append(clr)
        pPr.append(el)


def textbox(slide, x, y, w, h, lines, **kw):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    anchor = kw.pop('anchor', MSO_ANCHOR.TOP)
    tf.vertical_anchor = anchor
    fill_tf(tf, lines, **kw)
    return tb


def box(slide, x, y, w, h, lines, fill=LIGHT, line=None, radius=0.14,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE, pad=0.06, anchor=MSO_ANCHOR.MIDDLE,
        **kw):
    sh = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sh.adjustments[0] = radius
        except (IndexError, ValueError):
            pass
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.margin_left = tf.margin_right = Inches(pad)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    kw.setdefault('align', PP_ALIGN.CENTER)
    fill_tf(tf, lines, **kw)
    return sh


def line(slide, x1, y1, x2, y2, color=GRAY, width=0.75, arrow=False, dash=None):
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                    Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cn.line.color.rgb = color
    cn.line.width = Pt(width)
    ln = cn.line._get_or_add_ln()
    if arrow:
        ln.append(ln.makeelement(qn('a:tailEnd'),
                                 {'type': 'triangle', 'w': 'sm', 'len': 'sm'}))
    if dash:
        ln.append(ln.makeelement(qn('a:prstDash'), {'val': dash}))
    return cn


def rule(slide, x, y, w, color=ICON, width=1.25):
    """Линия-подчёркивание под заголовком блока."""
    return line(slide, x, y, x + w, y, color=color, width=width)


def block_title(slide, x, y, w, text, size=12, color=TEXT, ruled=True):
    tb = textbox(slide, x, y, w, 0.24, [text], size=size, bold=True, color=color)
    if ruled:
        rule(slide, x, y + 0.26, w)
    return tb


def chart_caption(slide, x, y, w, what, dim, size=11):
    """Название графика: «что, где, когда» чёрным + размерность серым."""
    return textbox(slide, x, y, w, 0.42,
                   [[(what + ', ', {'bold': True, 'size': size}),
                     (dim, {'bold': True, 'size': size, 'color': GRAY})]],
                   size=size)


def connect_tree(slide, px, py, cx, children_y, color=GRAY, spine_frac=0.45):
    """Соединяет узел-родителя (правый край px, центр py) с детьми (левый край cx)."""
    sx = px + (cx - px) * spine_frac
    line(slide, px, py, sx, py, color=color)
    if len(children_y) > 1:
        line(slide, sx, min(children_y), sx, max(children_y), color=color)
    for cy in children_y:
        line(slide, sx, cy, cx, cy, color=color, arrow=True)


# ------------------------------------------------------------------- слайды
def base_slide(prs, section, title, source, footnote=None):
    """Базовый слайд: раздел сверху, заголовок-вывод, источник в подвале."""
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    ph = {p.placeholder_format.idx: p for p in slide.placeholders}

    t = ph[0]
    # 11.42" — ширина из макета: правее начинается логотип
    t.width, t.height = Inches(11.42), Inches(0.80)
    t.left, t.top = Inches(LEFT), Inches(0.40)
    tf = t.text_frame
    tf.word_wrap = True
    fill_tf(tf, [title], size=20, bold=True, color=TEXT, line_spacing=0.95)

    # плейсхолдер колонтитула python-pptx не переносит из макета — рисуем сами
    textbox(slide, LEFT, 0.16, 11.42, 0.20, [section], size=12, color=GRAY)
    fill_tf(ph[12].text_frame, ['Источник: ' + source], size=9, color=GRAY)
    if footnote:
        fill_tf(ph[1].text_frame, [footnote], size=9, color=GRAY)
    else:
        ph[1]._element.getparent().remove(ph[1]._element)
    return slide


def section_slide(prs, title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    for p in list(slide.placeholders):
        p._element.getparent().remove(p._element)
    box(slide, 0, 3.05, SLIDE_W, 0.55, [title], fill=None, size=24, bold=True,
        color=TEXT, align=PP_ALIGN.CENTER)
    line(slide, SLIDE_W / 2 - 1.25, 3.72, SLIDE_W / 2 + 1.25, 3.72,
         color=ICON, width=1.5)
    if subtitle:
        textbox(slide, 0, 3.90, SLIDE_W, 0.3, [subtitle], size=12, color=GRAY,
                align=PP_ALIGN.CENTER)
    return slide


def title_slide(prs, title, subtitle, date=None):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    ph = {p.placeholder_format.idx: p for p in slide.placeholders}
    fill_tf(ph[0].text_frame, [title], size=32, bold=True, color=WHITE,
            line_spacing=1.05)
    fill_tf(ph[1].text_frame, [subtitle], size=13, color=WHITE)
    if date:
        textbox(slide, 5.89, 6.85, 7.08, 0.27, [date], size=12, color=WHITE,
                align=PP_ALIGN.RIGHT)
    return slide


def drop_first_slide(prs):
    """Убирает демонстрационные слайды шаблона вместе с их частями пакета."""
    lst = prs.slides._sldIdLst
    for sld in list(lst):
        prs.part.drop_rel(sld.rId)
        lst.remove(sld)


# -------------------------------------------------------------------- таблицы
def table(slide, x, y, w, col_w, rows, header=None, row_h=0.30, head_h=0.32,
          size=10, head_fill=C1, head_color=WHITE, zebra=(WHITE, LIGHT),
          aligns=None, bolds=None, colors=None, head_size=10):
    """Таблица из фигур: полный контроль над оформлением.

    col_w — доли ширины (сумма 1.0); rows — список списков строк;
    aligns/bolds/colors — по колонкам.
    """
    widths = [w * f for f in col_w]
    aligns = aligns or [PP_ALIGN.LEFT] * len(col_w)
    bolds = bolds or [False] * len(col_w)
    colors = colors or [TEXT] * len(col_w)
    cy = y
    if header:
        cx = x
        for j, cell in enumerate(header):
            box(slide, cx, cy, widths[j], head_h, [cell], fill=head_fill,
                radius=0.0, shape=MSO_SHAPE.RECTANGLE, size=head_size, bold=True,
                color=head_color, align=aligns[j], pad=0.07)
            cx += widths[j]
        cy += head_h
    for i, row in enumerate(rows):
        cx = x
        h = row_h(i) if callable(row_h) else row_h
        for j, cell in enumerate(row):
            fill = zebra[i % 2] if zebra else None
            kw = {}
            if isinstance(cell, tuple):
                cell, kw = cell
            box(slide, cx, cy, widths[j], h, [cell], fill=fill, radius=0.0,
                shape=MSO_SHAPE.RECTANGLE, size=kw.get('size', size),
                bold=kw.get('bold', bolds[j]), color=kw.get('color', colors[j]),
                align=kw.get('align', aligns[j]), pad=0.07, line=None)
            cx += widths[j]
        line(slide, x, cy + h, x + w, cy + h, color=GRAY_LT, width=0.5)
        cy += h
    return cy


def kpi(slide, x, y, w, h, value, caption, fill=LIGHT, value_color=ACCENT,
        value_size=17, caption_size=10):
    sh = box(slide, x, y, w, h, [], fill=fill, radius=0.10)
    tf = sh.text_frame
    fill_tf(tf, [[(value, {'size': value_size, 'bold': True, 'color': value_color})],
                 [(caption, {'size': caption_size, 'color': TEXT})]],
            align=PP_ALIGN.CENTER, line_spacing=0.95)
    return sh


def legend_row(slide, x, y, items, w_item, size=9.5, marker=0.13, gap_y=0.23,
               per_row=None):
    """Ручная легенда: цветной квадрат + подпись."""
    per_row = per_row or len(items)
    for i, (color, label) in enumerate(items):
        r, c = divmod(i, per_row)
        cx = x + c * w_item
        cy = y + r * gap_y
        box(slide, cx, cy + 0.03, marker, marker, [], fill=color, radius=0.0,
            shape=MSO_SHAPE.RECTANGLE)
        textbox(slide, cx + marker + 0.07, cy, w_item - marker - 0.12, 0.20,
                [label], size=size)
