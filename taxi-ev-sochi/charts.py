"""Графики под корпоративные правила шаблона.

Правила из шаблона: вертикальные гистограммы с накоплением; значения от
большего (снизу) к меньшему (сверху); цвета от тёмного к светлому; белая
обводка сегментов; сумма над столбцом; название графика — что/где/когда/
в какой размерности; шрифт 10 пт, не полужирный.
"""
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import (XL_CHART_TYPE, XL_LABEL_POSITION, XL_MARKER_STYLE,
                             XL_TICK_LABEL_POSITION)
from pptx.oxml.ns import qn

from deck_lib import GAMMA, TEXT, WHITE, GRAY, GRAY_LT, FONT


def _font(obj, size=10, bold=False, color=TEXT):
    obj.font.size = Pt(size)
    obj.font.bold = bold
    obj.font.name = FONT
    obj.font.color.rgb = color


def _kill_gridlines(ax):
    ax.has_major_gridlines = False
    ax.has_minor_gridlines = False


def _axis_off(ax):
    ax.visible = False
    _kill_gridlines(ax)


def _style_cat_axis(ax, size=10):
    _kill_gridlines(ax)
    ax.format.line.color.rgb = GRAY
    ax.format.line.width = Pt(0.75)
    _font(ax.tick_labels, size=size)
    ax.tick_labels.font.color.rgb = TEXT


# разряды отделяются пробелом только начиная с тысяч: иначе формат
# оставляет ведущий пробел у трёхзначных чисел
RU_NUM = '[>=1000]#\\ ##0;[<=-1000]-#\\ ##0;0'
RU_SIGNED = '+0;−0'


def stacked_columns(slide, x, y, w, h, categories, series, colors=None,
                    gap=90, label_size=10, min_label=60,
                    white_labels_upto=2):
    """Гистограмма с накоплением. series = [(имя, [значения...]), ...] снизу вверх."""
    data = CategoryChartData()
    data.categories = categories
    for name, vals in series:
        data.add_series(name, vals)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED,
                                Inches(x), Inches(y), Inches(w), Inches(h), data)
    ch = gf.chart
    ch.has_title = False
    ch.has_legend = False
    colors = colors or GAMMA
    plot = ch.plots[0]
    plot.gap_width = gap
    plot.overlap = 100
    for i, s in enumerate(ch.series):
        c = colors[i % len(colors)]
        s.format.fill.solid()
        s.format.fill.fore_color.rgb = c
        s.format.line.color.rgb = WHITE
        s.format.line.width = Pt(0.75)
        s.has_data_labels = True
        dl = s.data_labels
        dl.number_format = RU_NUM
        dl.number_format_is_linked = False
        dl.show_value = True
        dl.position = XL_LABEL_POSITION.CENTER
        _font(dl, size=label_size,
              color=WHITE if i < white_labels_upto else TEXT)
        # маленькие сегменты не подписываем — иначе каша
        for j, pt in enumerate(s.points):
            if abs(series[i][1][j] or 0) < min_label:
                _hide_label(pt)
    _style_cat_axis(ch.category_axis)
    _axis_off(ch.value_axis)
    return ch


def _hide_label(point):
    dLbl = point.data_label._get_or_add_dLbl()
    for child in list(dLbl):
        if child.tag != qn('c:idx'):
            dLbl.remove(child)
    dLbl.append(dLbl.makeelement(qn('c:delete'), {'val': '1'}))


def clustered_columns(slide, x, y, w, h, categories, values, point_colors,
                      gap=60, label_size=10, y_max=None, label_color=TEXT,
                      num_fmt=RU_NUM):
    data = CategoryChartData()
    data.categories = categories
    data.add_series('v', values)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                                Inches(x), Inches(y), Inches(w), Inches(h), data)
    ch = gf.chart
    ch.has_title = False
    ch.has_legend = False
    ch.plots[0].gap_width = gap
    s = ch.series[0]
    s.format.line.color.rgb = WHITE
    s.format.line.width = Pt(0.75)
    for pt, c in zip(s.points, point_colors):
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = c
    s.has_data_labels = True
    dl = s.data_labels
    dl.number_format = num_fmt
    dl.number_format_is_linked = False
    dl.show_value = True
    dl.position = XL_LABEL_POSITION.OUTSIDE_END
    _font(dl, size=label_size, bold=True, color=label_color)
    _style_cat_axis(ch.category_axis)
    _axis_off(ch.value_axis)
    if y_max:
        ch.value_axis.maximum_scale = y_max
        ch.value_axis.minimum_scale = 0
    return ch


def bar_tornado(slide, x, y, w, h, categories, values, point_colors,
                label_size=10, gap=45):
    """Горизонтальные полосы — для анализа чувствительности."""
    data = CategoryChartData()
    data.categories = categories
    data.add_series('v', values)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED,
                                Inches(x), Inches(y), Inches(w), Inches(h), data)
    ch = gf.chart
    ch.has_title = False
    ch.has_legend = False
    ch.plots[0].gap_width = gap
    s = ch.series[0]
    s.format.line.color.rgb = WHITE
    s.format.line.width = Pt(0.75)
    for pt, c in zip(s.points, point_colors):
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = c
    s.has_data_labels = True
    dl = s.data_labels
    dl.number_format = RU_SIGNED
    dl.number_format_is_linked = False
    dl.show_value = True
    dl.position = XL_LABEL_POSITION.OUTSIDE_END
    _font(dl, size=label_size, bold=True)
    cat = ch.category_axis
    _kill_gridlines(cat)
    cat.format.line.fill.background()
    cat.tick_label_position = XL_TICK_LABEL_POSITION.LOW
    _font(cat.tick_labels, size=label_size)
    _axis_off(ch.value_axis)
    return ch


def lines(slide, x, y, w, h, categories, series, colors, widths=None,
          dashes=None, label_size=10):
    data = CategoryChartData()
    data.categories = categories
    for name, vals in series:
        data.add_series(name, vals)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.LINE,
                                Inches(x), Inches(y), Inches(w), Inches(h), data)
    ch = gf.chart
    ch.has_title = False
    ch.has_legend = False
    for i, s in enumerate(ch.series):
        s.format.line.color.rgb = colors[i]
        s.format.line.width = Pt((widths or [2.25] * len(colors))[i])
        s.smooth = False
        if dashes and dashes[i]:
            ln = s.format.line._get_or_add_ln()
            ln.append(ln.makeelement(qn('a:prstDash'), {'val': dashes[i]}))
        s.marker.style = XL_MARKER_STYLE.NONE
    _style_cat_axis(ch.category_axis, size=label_size)
    va = ch.value_axis
    va.has_major_gridlines = True
    va.major_gridlines.format.line.color.rgb = GRAY_LT
    va.major_gridlines.format.line.width = Pt(0.5)
    va.format.line.fill.background()
    _font(va.tick_labels, size=label_size, color=GRAY)
    return ch


def manual_bars(slide, x, baseline, w, height, values, labels, colors,
                value_labels, bar_w=0.9, label_size=11, cat_size=10,
                axis_color=GRAY):
    """Столбчатая диаграмма из фигур: полный контроль над типографикой подписей.

    Нужна там, где формат чисел важнее автоматики (десятичная запятая и т. п.).
    """
    from deck_lib import box, line, textbox
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    vmax = max(values)
    n = len(values)
    step = w / n
    line(slide, x, baseline, x + w, baseline, color=axis_color, width=0.75)
    for i, v in enumerate(values):
        cx = x + step * (i + 0.5)
        h = height * v / vmax
        box(slide, cx - bar_w / 2, baseline - h, bar_w, h, [], fill=colors[i],
            radius=0.0, shape=MSO_SHAPE.RECTANGLE)
        textbox(slide, cx - step / 2, baseline - h - 0.26, step, 0.22,
                [value_labels[i]], size=label_size, bold=True,
                align=PP_ALIGN.CENTER)
        textbox(slide, cx - step / 2, baseline + 0.09, step, 0.44, [labels[i]],
                size=cat_size, align=PP_ALIGN.CENTER, line_spacing=0.98)
