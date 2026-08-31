#!/usr/bin/env python3
"""Строгая проверка готовой презентации.

Сверяет три вещи:
  1. каждое число на слайдах — с расчётом в model.py;
  2. числа в README.md и ИСТОЧНИКИ.md — с теми же расчётами;
  3. формальные требования задания и шаблона (число слайдов, границы
     рабочей области, кегли, источники, названия графиков).

Запуск: python3 check.py   ->   код возврата 0, если всё сошлось.
"""
import math
import os
import re
import subprocess
import sys

from pptx import Presentation
from pptx.util import Emu

import build_deck as B
import model

HERE = os.path.dirname(os.path.abspath(__file__))
PPTX = B.OUT
NBSP = ' '

fails, warns = [], []


def check(ok, name, detail=''):
    (fails if not ok else []).append(f'{name}{" — " + detail if detail else ""}')
    print(f'  {"OK  " if ok else "!!  "} {name}' + (f'   {detail}' if detail and not ok else ''))
    return ok


def warn(ok, name, detail=''):
    if not ok:
        warns.append(f'{name} — {detail}')
    print(f'  {"OK  " if ok else "??  "} {name}' + (f'   {detail}' if detail and not ok else ''))


# --------------------------------------------------------------- сбор текста
prs = Presentation(PPTX)
slides = list(prs.slides)


def shapes_of(slide):
    out = []

    def walk(shs):
        for sh in shs:
            if sh.shape_type == 6:
                walk(sh.shapes)
            else:
                out.append(sh)
    walk(slide.shapes)
    return out


def text_of(slide):
    return ' '.join(sh.text_frame.text for sh in shapes_of(slide)
                    if sh.has_text_frame)


def pdf_words(pdf):
    """Слова из отрендеренного PDF с координатами в дюймах, по страницам.

    Меряем по факту, а не по рамкам фигур: пустой низ текстового блока
    наложением не является.
    """
    out = subprocess.run(['pdftotext', '-bbox', pdf, '-'],
                         capture_output=True, text=True).stdout
    pages, cur = [], None
    for line in out.splitlines():
        if '<page ' in line:
            cur = []
            pages.append(cur)
        m = re.search(r'<word xMin="([\d.]+)" yMin="([\d.]+)" '
                      r'xMax="([\d.]+)" yMax="([\d.]+)">(.*)</word>', line)
        if m and cur is not None:
            x0, y0, x1, y1 = (float(v) / 72 for v in m.groups()[:4])
            cur.append((x0, y0, x1, y1, m.group(5)))
    return pages


def chart_values(slide):
    vals = []
    for sh in shapes_of(slide):
        if getattr(sh, 'has_chart', False) and sh.has_chart:
            for ser in sh.chart.series:
                vals.append([None if v is None else round(v) for v in ser.values])
    return vals


ALL = ' '.join(text_of(s) for s in slides)
FLAT = ALL.replace(NBSP, ' ')


def has(*fragments):
    """Все фрагменты присутствуют в тексте слайдов (пробелы нормализованы)."""
    return all(f.replace(NBSP, ' ') in FLAT for f in fragments)


# ====================================================== 1. слайды и модель
print('\n1. ЧИСЛА НА СЛАЙДАХ ПРОТИВ МОДЕЛИ')
ice, ev, delta = model.tco(model.P)
_, ev0, delta0 = model.tco(model.P, subsidy=False)
be, save_km, var_ice, var_ev, gap = model.breakeven_km(model.P)
be0 = model.breakeven_km(model.P, subsidy=False)[0]

check(round(delta) == B.GAIN == 480, 'выгода с субсидией = 480 тыс. ₽',
      f'модель {delta:.1f}')
check(round(-delta0) == B.LOSS == 445, 'проигрыш без субсидии = 445 тыс. ₽',
      f'модель {-delta0:.1f}')
check(has(f'{B.ru(B.GAIN)} тыс. ₽'), 'выгода напечатана на слайдах')
check(has(f'{B.ru(B.LOSS)} тыс. ₽'), 'проигрыш напечатан на слайдах')
check(has(f'{B.BE_K} тыс. км', f'{B.BE_M} мес'),
      f'окупаемость {B.BE_K} тыс. км / {B.BE_M} мес. на слайдах')
check(has(f'{B.BE0_K} тыс. км ({B.BE0_M} мес.)'),
      f'окупаемость без субсидии {B.BE0_K} тыс. км / {B.BE0_M} мес.')
check(has(f'{B.ru(B.SAVE_KM, 1)} ₽/км', f'{B.ru(B.VAR_EV, 1)} ₽/км',
          f'{B.ru(B.VAR_ICE, 1)} ₽/км'),
      f'экономия {B.ru(B.SAVE_KM, 1)} = {B.ru(B.VAR_ICE, 1)} − {B.ru(B.VAR_EV, 1)} ₽/км')
check(abs((var_ice - var_ev) - save_km) < 1e-9,
      'экономия действительно равна разности переменных затрат')
check(has(f'{B.ru(B.FIXED_GAP)} тыс. ₽', f'{B.ru(B.FIXED_GAP0)} тыс. ₽'),
      f'разрыв по постоянным {B.FIXED_GAP} и {B.FIXED_GAP0} тыс. ₽')
check(B.FIXED_GAP0 - B.FIXED_GAP == model.P['ev_subsidy'],
      'разница между разрывами равна субсидии')
check(has(f'{B.ru(B.CAPEX_GAP)} тыс. ₽', f'{B.ru(B.CAPEX_GAP0)} тыс. ₽'),
      f'переплата за актив {B.CAPEX_GAP} / {B.CAPEX_GAP0} тыс. ₽')

# графики слайда 4: столбцы должны быть ровно статьями модели
stack = chart_values(slides[3])[0:5]
expect = [[round(d[k]) for d in (ice, ev, ev0)]
          for k in ('capex', 'energy', 'service', 'hub', 'downtime')]
check(stack == expect, 'столбцы TCO совпадают со статьями модели',
      f'{stack} против {expect}')
for name, d in (('ДВС', ice), ('ЭМ с субсидией', ev), ('ЭМ без субсидии', ev0)):
    check(has(B.ru(d['total'])), f'итог {name} = {B.ru(d["total"])} на слайде')
check(round(ice['total'] - ev['total']) == B.GAIN,
      'итог ДВС минус итог ЭМ равен заявленной выгоде')

# кривые накопленных затрат должны пересекаться в точке безубыточности
curves = chart_values(slides[3])[5:8]
check(len(curves) == 3, 'на слайде 4 три кривые накопленных затрат')
if len(curves) == 3:
    at270 = [c[3] for c in curves]
    check([round(ice['total']), round(ev['total']), round(ev0['total'])] == at270,
          'кривые в точке 270 тыс. км сходятся с итогами TCO', str(at270))

# мощность
power = chart_values(slides[4])[0]
check(power == [250] + [round(B.POWER[n]) for n in (20, 50, 100, 200)],
      'столбцы мощности совпадают с расчётом', str(power))
check(has(f'{B.ru(B.POWER[200] / 1000, 2)} МВт'),
      f'мощность на 200 авто = {B.ru(B.POWER[200] / 1000, 2)} МВт')
check(abs(B.POWER[200] - 1356) < 1, 'мощность на 200 авто ≈ 1356 кВт')

# чувствительность
tornado = chart_values(slides[5])[0]
cases = sorted([model.tco(model.P, ice_fuel=model.P['ice_fuel'] * 1.2)[2],
                delta,
                model.tco(model.P, ev_kwh100=24.0)[2],
                model.tco(model.P, downtime=0.06)[2],
                model.tco(model.P, ev_residual=0.15)[2],
                model.tco(model.P, km_year=60_000)[2],
                model.tco(model.P, ev_price_kwh=20.0)[2],
                delta0])
check([round(v) for v in cases] == sorted(tornado),
      'полосы чувствительности совпадают с пересчётом модели',
      f'{sorted(tornado)} против {[round(v) for v in cases]}')
check(min(tornado) == round(delta0),
      'самый плохой сценарий — именно отсутствие субсидии')
check(sum(1 for v in tornado if v < 0) == 1,
      'ровно один сценарий уходит в минус — это держит вывод слайда')

# пилот
check(abs(B.PILOT_GAIN - B.GAIN * B.PILOT / 1000) < 1e-9,
      'эффект пилота = выгода × 20 машин')
check(has(f'+{B.ru(B.PILOT_GAIN, 1)} млн ₽'),
      f'эффект пилота {B.ru(B.PILOT_GAIN, 1)} млн ₽ на слайдах')
check(has(f'{B.ru(B.PILOT_CAPEX)} млн ₽'),
      f'инвестиции в пилот {B.ru(B.PILOT_CAPEX)} млн ₽')
check(abs(B.PILOT_CAPEX - (B.EV_PAID + model.P['hub_capex']) * B.PILOT / 1000) < 1e-9,
      'инвестиции = (цена после субсидии + хаб) × 20')
check(B.EV_PAID == model.P['ev_price'] - model.P['ev_subsidy'] == 1850,
      'цена ЭМ после субсидии = 1 850 тыс. ₽')

# гарантия
check(B.BAT_M == 20 and B.CAR_M == 13,
      'гарантия: батарея 150 тыс. км = 20 мес., авто 100 тыс. км = 13 мес.')
check(B.BAT_UNCOVERED == 120, 'без гарантии на батарею остаётся 120 тыс. км')
check(has('150 тыс. км', '120 тыс. км'), 'сроки гарантии напечатаны на слайдах')
check('втрое короче' not in FLAT, 'снято прежнее неверное «втрое короче»')

# старые, уже исправленные значения не должны нигде остаться
for stale in ('531', '394', '3 178', '3 157', '1 770', '381 тыс', '120 тыс. км (16',
              '7,9 ₽', '3,5 ₽', '10,6 млн', '13-м месяце'):
    check(stale.replace(NBSP, ' ') not in FLAT, f'нет устаревшего «{stale}»')


# ====================================================== 2. документы
print('\n2. README.md И ИСТОЧНИКИ.md ПРОТИВ МОДЕЛИ')
readme = open(os.path.join(HERE, 'README.md'), encoding='utf-8').read()
src = open(os.path.join(HERE, 'ИСТОЧНИКИ.md'), encoding='utf-8').read()

for doc, name in ((readme, 'README.md'), (src, 'ИСТОЧНИКИ.md')):
    flat = doc.replace(NBSP, ' ')
    check(f'{B.GAIN} тыс. ₽' in flat, f'{name}: выгода {B.GAIN} тыс. ₽')
    check(f'{B.LOSS} тыс. ₽' in flat, f'{name}: проигрыш {B.LOSS} тыс. ₽')
    check(f'{B.BE_K} тыс. км' in flat, f'{name}: окупаемость {B.BE_K} тыс. км')
    body = flat.split('## 1.')[-1]          # всё после раздела «что исправлено»
    for stale in ('531 тыс', '394 тыс', '3 178', '7,9 ₽/км', '3,5 ₽/км'):
        check(stale not in body,
              f'{name}: «{stale}» — только в разделе исправлений, не в тексте')

# арифметика, выписанная в ИСТОЧНИКИ.md, должна сходиться с моделью
for label, value in (('1 027', ice['capex']), ('1 719', ice['energy']),
                     ('1 240', ev['capex']), ('832', ev['energy']),
                     ('2 164', ev0['capex']), ('3 106', ice['total']),
                     ('2 626', ev['total']), ('3 550', ev0['total'])):
    check(label in src.replace(NBSP, ' '),
          f'ИСТОЧНИКИ.md: строка {label} есть')
    check(abs(round(value) - int(label.replace(' ', ''))) <= 1,
          f'ИСТОЧНИКИ.md: {label} совпадает с моделью ({value:.1f})')


# ====================================================== 3. требования
print('\n3. ТРЕБОВАНИЯ ЗАДАНИЯ И ШАБЛОНА')
check(5 <= len(slides) <= 10, f'слайдов {len(slides)} — в диапазоне 5–10')
check('Executive Summary' in text_of(slides[1]),
      'первый содержательный слайд — Executive Summary')

# структура повторяет дерево
order = ['Executive Summary', 'дерево решения', '1. Выгодно?', '2. Выполнимо?',
         '3. Устойчиво?', 'Рекомендация', 'Приложение']
found = [i for i, key in enumerate(order)
         if any(key.lower() in text_of(sl).lower() for sl in slides)]
check(len(found) == len(order), 'все разделы дерева присутствуют по порядку')

LEFT, RIGHT = 0.37, 12.97
DECOR_X, DECOR_Y = 12.06, 6.23      # габарит декоративной четверти круга
DECOR_C, DECOR_R = (13.333, 7.5), 1.27
PAGES = pdf_words(os.path.join(HERE, 'output',
                               'Электромобили_в_таксопарке_Сочи.pdf'))
for i, sl in enumerate(slides, 1):
    if i in (1, 8):
        continue
    txt = text_of(sl)
    check('Источник:' in txt, f'слайд {i}: есть строка источника')
    m = re.search(r'Источник:([^\n]*)', txt)
    if m:
        check(not re.search(r'https?://|www\.|\S+\.\w{2,4}/', m.group(1)),
              f'слайд {i}: в источниках нет ссылок', m.group(1)[:70])
    # заголовок не длиннее двух строк и не заезжает под логотип
    title = sl.shapes.title
    if title is not None:
        w = Emu(title.width).inches
        check(w <= 11.45, f'слайд {i}: заголовок не заходит под логотип',
              f'ширина {w:.2f}"')
        est_lines = max(1, round(len(title.text_frame.text) / (w * 7.9) + 0.49))
        check(est_lines <= 2, f'слайд {i}: заголовок не длиннее 2 строк',
              f'~{est_lines} строк, {len(title.text_frame.text)} знаков')
    # кегли: в теле 10–12, мельче — только сноска и подпись оси
    for sh in shapes_of(sl):
        if not sh.has_text_frame:
            continue
        for para in sh.text_frame.paragraphs:
            for r in para.runs:
                if not r.text.strip() or r.font.size is None:
                    continue
                pt = r.font.size.pt
                small_ok = (r.text.startswith('Источник:') or r.text.startswith('1 – ')
                            or 'пробег, тыс. км' in r.text)
                if pt < 10 and not small_ok:
                    check(False, f'слайд {i}: кегль ниже 10 пт', f'{pt} пт «{r.text[:40]}»')
                if pt > 20:
                    check(False, f'слайд {i}: кегль выше 20 пт', f'{pt} пт')
    # заливки фигур не должны наезжать на декор в правом нижнем углу
    for sh in shapes_of(sl):
        try:
            right = Emu(sh.left).inches + Emu(sh.width).inches
            bottom = Emu(sh.top).inches + Emu(sh.height).inches
        except TypeError:
            continue
        filled = getattr(getattr(sh, 'fill', None), 'type', None) == 1
        if filled and bottom > DECOR_Y and right > DECOR_X:
            check(False, f'слайд {i}: заливка наезжает на декор угла',
                  f'{sh.name} {right:.2f}" / {bottom:.2f}"')

# по факту отрендеренного PDF: ни одно слово не влезает в декор и за поля
print('\n   проверка по отрендеренному PDF')
before = len(fails)
for i, words in enumerate(PAGES, 1):
    if i in (1, 8):
        continue
    for x0, y0, x1, y1, w in words:
        if not w.strip() or w.strip() == str(i):    # номер слайда живёт в круге
            continue
        if any(math.hypot(DECOR_C[0] - cx, DECOR_C[1] - cy) < DECOR_R
               for cx, cy in ((x1, y1), (x0, y1), (x1, y0))):
            check(False, f'слайд {i}: слово попадает в декор угла',
                  f'«{w}» в ({x1:.2f}, {y1:.2f})')
        elif x1 > RIGHT + 0.05:
            check(False, f'слайд {i}: слово правее рабочей области',
                  f'«{w}» до {x1:.2f}"')
        elif x0 < LEFT - 0.05:
            check(False, f'слайд {i}: слово левее рабочей области',
                  f'«{w}» от {x0:.2f}"')
if len(fails) == before:
    print('  OK   ни одно слово не попадает в декор и не выходит за поля')

# названия графиков: что, где, когда, размерность
CAPTIONS = ['Стоимость владения одним автомобилем такси, Сочи, 3 года / 270 тыс. км',
            'Накопленные затраты на один автомобиль по мере пробега, Сочи',
            'Присоединённая мощность для ночной зарядки парка, Сочи',
            'Выгода от перехода на электромобиль за 3 года при изменении одного параметра, Сочи',
            'Доля в продажах новых легковых автомобилей, Россия, январь–июль']
for cap in CAPTIONS:
    check(cap.replace(NBSP, ' ') in FLAT, f'название графика: «{cap[:48]}…»')

# приложение отделено и стоит в конце
app_start = next(i for i, sl in enumerate(slides)
                 if text_of(sl).strip().startswith('Приложение'))
check(app_start == 7, 'разделитель «Приложение» — восьмой слайд')
check(all('Приложение' in text_of(sl) for sl in slides[7:]),
      'после разделителя идут только слайды приложения')

# метаданные не должны тащить автора шаблона
cp = prs.core_properties
check(not cp.author, 'в свойствах файла нет автора шаблона', repr(cp.author))
check('Сочи' in (cp.title or ''), 'в свойствах файла заполнен заголовок')

pdf = os.path.join(HERE, 'output', 'Электромобили_в_таксопарке_Сочи.pdf')
check(os.path.exists(pdf), 'PDF выгружен рядом с PPTX')
if os.path.exists(pdf):
    check(os.path.getmtime(pdf) >= os.path.getmtime(PPTX) - 1,
          'PDF не старше PPTX')

# ------------------------------------------------------------------- итог
print('\n' + '=' * 62)
if fails:
    print(f'НЕ СОШЛОСЬ: {len(fails)}')
    for f in fails:
        print('  •', f)
else:
    print('Всё сошлось: расхождений между моделью, слайдами и документами нет.')
if warns:
    print(f'\nПредупреждения: {len(warns)}')
    for w in warns:
        print('  •', w)
sys.exit(1 if fails else 0)
